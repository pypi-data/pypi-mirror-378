#!/usr/bin/env python3
"""
MCP server for PowerPoint file manipulation
Provides tools to convert PowerPoint files to/from JSON using Pydantic models
"""

import asyncio
import json
import base64
import io
from pathlib import Path
from typing import Optional, Dict, Any, List
from datetime import datetime

from fastmcp import FastMCP
from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.dml import MSO_FILL_TYPE, MSO_LINE_DASH_STYLE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image

from .powerpoint_models import (
    Presentation, Slide, Shape, TextFrame, Paragraph, TextRun, 
    FontInfo, ParagraphFormat, TableInfo, TableCell, ImageInfo,
    FillFormat, LineFormat, ShadowFormat, DocumentProperties,
    ShapeType
)

# Create the FastMCP app
mcp = FastMCP("PowerPoint Tools")

def rgb_to_hex(color_obj) -> Optional[str]:
    """Convert color object to hex string"""
    if color_obj is None:
        return None
    try:
        # Handle different color object types
        if hasattr(color_obj, 'rgb'):
            return f"#{color_obj.rgb:06X}"
        elif hasattr(color_obj, 'color') and hasattr(color_obj.color, 'rgb'):
            return f"#{color_obj.color.rgb:06X}"
        return None
    except:
        return None

def hex_to_rgb(hex_color: str) -> Optional[RGBColor]:
    """Convert hex color string to RGBColor"""
    if not hex_color or not hex_color.startswith('#'):
        return None
    try:
        rgb_int = int(hex_color[1:], 16)
        return RGBColor((rgb_int >> 16) & 0xFF, (rgb_int >> 8) & 0xFF, rgb_int & 0xFF)
    except:
        return None

def extract_font_info(font) -> Optional[FontInfo]:
    """Extract font information from python-pptx font object"""
    if not font:
        return None
    
    # Convert EMU to points for font size
    font_size = None
    if hasattr(font, 'size') and font.size is not None:
        font_size = int(font.size / 12700)  # EMU to points conversion
    
    # Extract font color properly
    font_color = None
    if hasattr(font, 'color') and font.color is not None:
        try:
            if hasattr(font.color, 'rgb') and font.color.rgb is not None:
                rgb_val = font.color.rgb
                if hasattr(rgb_val, '_rgb_val'):
                    font_color = f"#{rgb_val._rgb_val:06X}"
                else:
                    try:
                        font_color = f"#{int(rgb_val):06X}"
                    except:
                        rgb_str = str(rgb_val)
                        if len(rgb_str) == 6:
                            font_color = f"#{rgb_str.upper()}"
            else:
                # If RGB can't be accessed, it might be an automatic/theme color
                # For most text, automatic color defaults to black
                font_color = "#000000"
        except:
            # Default to black for text when color extraction fails
            font_color = "#000000"
    
    return FontInfo(
        name=getattr(font, 'name', None),
        size=font_size,
        bold=getattr(font, 'bold', None),
        italic=getattr(font, 'italic', None),
        underline=getattr(font, 'underline', None),
        color_rgb=font_color
    )

def extract_paragraph_format(paragraph) -> Optional[ParagraphFormat]:
    """Extract paragraph formatting"""
    if not hasattr(paragraph, 'format'):
        return None
    
    fmt = paragraph.format
    alignment_map = {
        PP_ALIGN.LEFT: "LEFT",
        PP_ALIGN.CENTER: "CENTER", 
        PP_ALIGN.RIGHT: "RIGHT",
        PP_ALIGN.JUSTIFY: "JUSTIFY"
    }
    
    return ParagraphFormat(
        alignment=alignment_map.get(getattr(fmt, 'alignment', None)),
        space_before=getattr(fmt, 'space_before', None),
        space_after=getattr(fmt, 'space_after', None),
        line_spacing=getattr(fmt, 'line_spacing', None),
        level=getattr(fmt, 'level', None)
    )

def extract_text_frame(text_frame) -> Optional[TextFrame]:
    """Extract text frame information"""
    if not text_frame:
        return None
    
    paragraphs = []
    for p in text_frame.paragraphs:
        runs = []
        for run in p.runs:
            runs.append(TextRun(
                text=run.text,
                font=extract_font_info(run.font)
            ))
        
        paragraphs.append(Paragraph(
            text=p.text,
            runs=runs,
            format=extract_paragraph_format(p)
        ))
    
    return TextFrame(
        text=text_frame.text,
        paragraphs=paragraphs,
        margin_left=getattr(text_frame, 'margin_left', None),
        margin_right=getattr(text_frame, 'margin_right', None),
        margin_top=getattr(text_frame, 'margin_top', None),
        margin_bottom=getattr(text_frame, 'margin_bottom', None),
        word_wrap=getattr(text_frame, 'word_wrap', None)
    )

def extract_table_info(table) -> Optional[TableInfo]:
    """Extract table information"""
    if not table:
        return None
    
    cells = []
    for row in table.rows:
        cell_row = []
        for cell in row.cells:
            cell_row.append(TableCell(
                text=cell.text,
                text_frame=extract_text_frame(cell.text_frame)
            ))
        cells.append(cell_row)
    
    return TableInfo(
        rows=len(table.rows),
        columns=len(table.columns),
        cells=cells
    )

def extract_image_info(shape) -> Optional[ImageInfo]:
    """Extract image information and convert to base64"""
    try:
        image = shape.image
        image_data = base64.b64encode(image.blob).decode('utf-8')
        
        return ImageInfo(
            filename=getattr(image, 'filename', None),
            image_data=image_data,
            content_type=image.content_type,
            crop_left=getattr(shape, 'crop_left', None),
            crop_top=getattr(shape, 'crop_top', None),
            crop_right=getattr(shape, 'crop_right', None),
            crop_bottom=getattr(shape, 'crop_bottom', None)
        )
    except:
        return None

def determine_shape_type(shape) -> ShapeType:
    """Determine the shape type from MSO_SHAPE_TYPE"""
    type_mapping = {
        MSO_SHAPE_TYPE.AUTO_SHAPE: ShapeType.AUTO_SHAPE,
        MSO_SHAPE_TYPE.TEXT_BOX: ShapeType.TEXT_BOX,
        MSO_SHAPE_TYPE.PICTURE: ShapeType.PICTURE,
        MSO_SHAPE_TYPE.TABLE: ShapeType.TABLE,
        MSO_SHAPE_TYPE.PLACEHOLDER: ShapeType.PLACEHOLDER,
        MSO_SHAPE_TYPE.GROUP: ShapeType.GROUP,
        MSO_SHAPE_TYPE.CHART: ShapeType.CHART,
        MSO_SHAPE_TYPE.FREEFORM: ShapeType.FREEFORM
    }
    return type_mapping.get(shape.shape_type, ShapeType.OTHER)

def extract_fill_format(shape) -> Optional[FillFormat]:
    """Extract fill formatting"""
    try:
        if not hasattr(shape, 'fill'):
            return None
        
        fill = shape.fill
        fill_type_map = {
            MSO_FILL_TYPE.SOLID: "SOLID",
            MSO_FILL_TYPE.GRADIENT: "GRADIENT", 
            MSO_FILL_TYPE.PICTURE: "PICTURE",
            MSO_FILL_TYPE.PATTERNED: "PATTERNED",
            MSO_FILL_TYPE.TEXTURED: "TEXTURED",
            MSO_FILL_TYPE.BACKGROUND: "BACKGROUND"
        }
        
        # Extract fill colors properly using string method (works around RGBColor formatting issue)
        fore_color = None
        back_color = None
        fill_type = getattr(fill, 'type', None)
        
        # Only try to access colors for appropriate fill types to avoid exceptions
        if fill_type == MSO_FILL_TYPE.SOLID:
            try:
                if hasattr(fill, 'fore_color') and fill.fore_color is not None:
                    if hasattr(fill.fore_color, 'rgb') and fill.fore_color.rgb is not None:
                        rgb_str = str(fill.fore_color.rgb)
                        if len(rgb_str) == 6:
                            fore_color = f"#{rgb_str.upper()}"
            except:
                pass
        
        elif fill_type == MSO_FILL_TYPE.GRADIENT or fill_type == MSO_FILL_TYPE.PATTERNED:
            try:
                if hasattr(fill, 'fore_color') and fill.fore_color is not None:
                    if hasattr(fill.fore_color, 'rgb') and fill.fore_color.rgb is not None:
                        rgb_str = str(fill.fore_color.rgb)
                        if len(rgb_str) == 6:
                            fore_color = f"#{rgb_str.upper()}"
                            
                if hasattr(fill, 'back_color') and fill.back_color is not None:
                    if hasattr(fill.back_color, 'rgb') and fill.back_color.rgb is not None:
                        rgb_str = str(fill.back_color.rgb)
                        if len(rgb_str) == 6:
                            back_color = f"#{rgb_str.upper()}"
            except:
                pass
        
        return FillFormat(
            fill_type=fill_type_map.get(getattr(fill, 'type', None)),
            fore_color=fore_color,
            back_color=back_color,
            transparency=getattr(fill, 'transparency', None)
        )
    except:
        return None

def extract_line_format(shape) -> Optional[LineFormat]:
    """Extract line formatting"""
    try:
        if not hasattr(shape, 'line'):
            return None
        
        line = shape.line
        
        # Extract line color properly using string method 
        line_color = None
        if hasattr(line, 'color') and line.color is not None:
            try:
                if hasattr(line.color, 'rgb') and line.color.rgb is not None:
                    rgb_str = str(line.color.rgb)
                    if len(rgb_str) == 6:
                        line_color = f"#{rgb_str.upper()}"
            except:
                # If color RGB can't be accessed, mark as no line
                line_color = "NO_LINE"
        
        return LineFormat(
            color=line_color,
            width=getattr(line, 'width', None),
            transparency=getattr(line, 'transparency', None)
        )
    except:
        return None

def extract_shape(shape) -> Shape:
    """Extract complete shape information"""
    shape_data = Shape(
        shape_id=shape.shape_id,
        name=shape.name,
        shape_type=determine_shape_type(shape),
        left=shape.left,
        top=shape.top,
        width=shape.width,
        height=shape.height,
        rotation=getattr(shape, 'rotation', None)
    )
    
    # Extract text frame if present
    if hasattr(shape, 'text_frame') and shape.text_frame:
        shape_data.text_frame = extract_text_frame(shape.text_frame)
    
    # Extract table if present
    if hasattr(shape, 'table') and shape.table:
        shape_data.table = extract_table_info(shape.table)
    
    # Extract image if present
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        shape_data.image = extract_image_info(shape)
    
    # Extract formatting
    shape_data.fill = extract_fill_format(shape)
    shape_data.line = extract_line_format(shape)
    
    # Handle grouped shapes
    if hasattr(shape, 'shapes'):
        shape_data.shapes = [extract_shape(s) for s in shape.shapes]
    
    return shape_data

@mcp.tool()
def pptx_to_json(file_path: str) -> str:
    """
    Convert a PowerPoint file to JSON format using Pydantic models.
    
    Args:
        file_path: Path to the .pptx file
        
    Returns:
        JSON string representing the complete presentation structure
    """
    try:
        # Load the presentation
        prs = PptxPresentation(file_path)
        
        # Extract slides
        slides = []
        for slide_idx, slide in enumerate(prs.slides):
            shapes = [extract_shape(shape) for shape in slide.shapes]
            
            # Extract notes safely
            notes_text = None
            try:
                if hasattr(slide, 'notes_slide') and slide.notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text if hasattr(slide.notes_slide, 'notes_text_frame') else None
            except:
                notes_text = None
            
            slide_data = Slide(
                slide_number=slide_idx + 1,
                slide_id=getattr(slide, 'slide_id', None),
                shapes=shapes,
                notes=notes_text
            )
            slides.append(slide_data)
        
        # Extract document properties
        core_props = None
        if hasattr(prs, 'core_properties'):
            cp = prs.core_properties
            core_props = DocumentProperties(
                title=getattr(cp, 'title', None),
                author=getattr(cp, 'author', None),
                subject=getattr(cp, 'subject', None),
                keywords=getattr(cp, 'keywords', None),
                comments=getattr(cp, 'comments', None),
                created=getattr(cp, 'created', None).isoformat() if getattr(cp, 'created', None) else None,
                modified=getattr(cp, 'modified', None).isoformat() if getattr(cp, 'modified', None) else None,
                last_modified_by=getattr(cp, 'last_modified_by', None)
            )
        
        # Create the presentation model
        presentation = Presentation(
            slide_width=prs.slide_width,
            slide_height=prs.slide_height,
            slides=slides,
            core_properties=core_props
        )
        
        return presentation.model_dump_json(indent=2)
        
    except Exception as e:
        return json.dumps({"error": f"Failed to process PowerPoint file: {str(e)}"})

@mcp.tool() 
def json_to_pptx(json_data: str, output_path: str) -> str:
    """
    Create a PowerPoint file from JSON data.
    
    Args:
        json_data: JSON string representing the presentation structure
        output_path: Path where the .pptx file should be saved
        
    Returns:
        Success message or error details
    """
    try:
        # Parse the JSON data
        presentation_dict = json.loads(json_data)
        presentation = Presentation(**presentation_dict)
        
        # Create new presentation
        prs = PptxPresentation()
        
        # Set slide dimensions
        prs.slide_width = presentation.slide_width
        prs.slide_height = presentation.slide_height
        
        # Clear default slides
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
        
        # Add slides
        for slide_data in presentation.slides:
            # Add blank slide
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add shapes
            for shape_data in slide_data.shapes:
                add_shape_to_slide(slide, shape_data)
        
        # Set document properties
        if presentation.core_properties:
            cp = prs.core_properties
            if presentation.core_properties.title:
                cp.title = presentation.core_properties.title
            if presentation.core_properties.author:
                cp.author = presentation.core_properties.author
            if presentation.core_properties.subject:
                cp.subject = presentation.core_properties.subject
        
        # Save the presentation
        prs.save(output_path)
        
        return f"Successfully created PowerPoint file: {output_path}"
        
    except Exception as e:
        return f"Failed to create PowerPoint file: {str(e)}"

def add_shape_to_slide(slide, shape_data: Shape):
    """Add a shape to a slide based on shape data"""
    try:
        created_shape = None
        
        if shape_data.shape_type == ShapeType.TEXT_BOX:
            # Add text box
            created_shape = slide.shapes.add_textbox(
                shape_data.left, shape_data.top, 
                shape_data.width, shape_data.height
            )
            if shape_data.text_frame:
                created_shape.text = shape_data.text_frame.text
                apply_text_formatting(created_shape.text_frame, shape_data.text_frame)
        
        elif shape_data.shape_type == ShapeType.AUTO_SHAPE:
            # Add auto shape (rectangle as default)
            from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
            created_shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                shape_data.left, shape_data.top,
                shape_data.width, shape_data.height
            )
            if shape_data.text_frame:
                created_shape.text = shape_data.text_frame.text
                apply_text_formatting(created_shape.text_frame, shape_data.text_frame)
        
        elif shape_data.shape_type == ShapeType.PICTURE and shape_data.image:
            # Add picture
            if shape_data.image.image_data:
                image_data = base64.b64decode(shape_data.image.image_data)
                image_stream = io.BytesIO(image_data)
                created_shape = slide.shapes.add_picture(
                    image_stream, shape_data.left, shape_data.top,
                    shape_data.width, shape_data.height
                )
        
        elif shape_data.shape_type == ShapeType.TABLE and shape_data.table:
            # Add table
            table_shape = slide.shapes.add_table(
                shape_data.table.rows, shape_data.table.columns,
                shape_data.left, shape_data.top, 
                shape_data.width, shape_data.height
            )
            table = table_shape.table
            created_shape = table_shape
            
            # Populate table cells
            if shape_data.table.cells:
                for row_idx, row in enumerate(shape_data.table.cells):
                    for col_idx, cell_data in enumerate(row):
                        if row_idx < len(table.rows) and col_idx < len(table.columns):
                            cell = table.cell(row_idx, col_idx)
                            cell.text = cell_data.text
        
        # Apply fill formatting to the created shape
        if created_shape and shape_data.fill:
            apply_fill_formatting(created_shape, shape_data.fill)
        
        # Apply line formatting to the created shape  
        if created_shape and shape_data.line:
            apply_line_formatting(created_shape, shape_data.line)
        
    except Exception as e:
        print(f"Warning: Failed to add shape {shape_data.name}: {str(e)}")

def apply_fill_formatting(shape, fill_data: FillFormat):
    """Apply fill formatting to a shape"""
    try:
        if not hasattr(shape, 'fill'):
            return
            
        fill = shape.fill
        
        if fill_data.fill_type == "SOLID" and fill_data.fore_color:
            fill.solid()
            rgb_color = hex_to_rgb(fill_data.fore_color)
            if rgb_color:
                fill.fore_color.rgb = rgb_color
        elif fill_data.fill_type == "BACKGROUND":
            # For background fills, explicitly set to no fill to prevent default gradient
            fill.background()
        
    except Exception as e:
        print(f"Warning: Failed to apply fill formatting: {str(e)}")

def apply_line_formatting(shape, line_data: LineFormat):
    """Apply line formatting to a shape"""
    try:
        if not hasattr(shape, 'line'):
            return
            
        line = shape.line
        
        if line_data.color == "NO_LINE" or line_data.color is None:
            # Explicitly set no line to prevent default blue borders
            try:
                # Set line fill to no fill
                line.fill.background()
                line.width = 0
            except:
                try:
                    # Alternative approach - set width to 0
                    line.width = 0
                except:
                    pass
        elif line_data.color:
            rgb_color = hex_to_rgb(line_data.color)
            if rgb_color:
                line.color.rgb = rgb_color
        
        if line_data.width is not None:
            line.width = line_data.width
            
    except Exception as e:
        print(f"Warning: Failed to apply line formatting: {str(e)}")

def apply_text_formatting(text_frame, text_frame_data: TextFrame):
    """Apply text formatting to a text frame"""
    try:
        if text_frame_data.paragraphs:
            # Clear existing paragraphs except the first
            try:
                while len(text_frame.paragraphs) > 1:
                    # Try different ways to delete paragraphs
                    if hasattr(text_frame.paragraphs, '_delete_paragraph'):
                        text_frame.paragraphs._delete_paragraph(text_frame.paragraphs[-1])
                    else:
                        # Alternative approach - clear all and rebuild
                        text_frame.clear()
                        break
            except:
                # If paragraph deletion fails, clear and rebuild
                text_frame.clear()
            
            # Set first paragraph or add new ones
            for i, para_data in enumerate(text_frame_data.paragraphs):
                if i == 0:
                    para = text_frame.paragraphs[0]
                else:
                    para = text_frame.add_paragraph()
                
                para.text = para_data.text
                
                # Apply paragraph formatting
                if para_data.format:
                    if para_data.format.alignment:
                        alignment_map = {
                            "LEFT": PP_ALIGN.LEFT,
                            "CENTER": PP_ALIGN.CENTER,
                            "RIGHT": PP_ALIGN.RIGHT,
                            "JUSTIFY": PP_ALIGN.JUSTIFY
                        }
                        if para_data.format.alignment in alignment_map:
                            para.alignment = alignment_map[para_data.format.alignment]
                
                # Apply run formatting
                if para_data.runs:
                    para.clear()
                    for run_data in para_data.runs:
                        run = para.add_run()
                        run.text = run_data.text
                        if run_data.font:
                            apply_font_formatting(run.font, run_data.font)
    
    except Exception as e:
        print(f"Warning: Failed to apply text formatting: {str(e)}")

def apply_font_formatting(font, font_data: FontInfo):
    """Apply font formatting"""
    try:
        if font_data.name:
            font.name = font_data.name
        if font_data.size is not None:
            font.size = Pt(font_data.size)  # font_data.size is already in points
        if font_data.bold is not None:
            font.bold = font_data.bold
        if font_data.italic is not None:
            font.italic = font_data.italic
        if font_data.color_rgb:
            rgb_color = hex_to_rgb(font_data.color_rgb)
            if rgb_color:
                font.color.rgb = rgb_color
            else:
                print(f"Warning: Could not convert font color {font_data.color_rgb} to RGB")
    except Exception as e:
        print(f"Warning: Failed to apply font formatting: {str(e)}")

def main():
    """Main entry point for the MCP server"""
    mcp.run()

if __name__ == "__main__":
    main()