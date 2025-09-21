"""
📄 extract_utils.py

Purpose:
    Contains file-type specific text extractors (DOCX, EML, MSG).

Key Features:
    - _extract_docx(): Parses Word documents and extracts text and tables.
    - _extract_eml(): Parses .eml email files using eml_parser.
    - _extract_msg(): Parses .msg Outlook files using extract_msg.
    - Each extractor also supports virtual tag detection via fts_core.

Usage:
    Called by `filetype_utils.py -> extract_text_from_file()` during indexing.
"""

import io
import re
import os
import sqlite3
import docx
import extract_msg
import eml_parser
import json
import fitz  # PyMuPDF
import pytesseract
import openpyxl



from pptx import Presentation
from ebooklib import epub
from bs4 import BeautifulSoup
from odf.opendocument import load
from odf.text import P
from PIL import Image, ExifTags
from datetime import datetime

from .config import DB_FILE


def _extract_docx(path):
    from .fts_core import extract_virtual_tags

    doc = docx.Document(path)

    # Extract paragraphs
    raw_paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    # Extract structured table info
    table_lines = []
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if len(cells) >= 2:
                key = cells[0]
                value = " ".join(cells[1:])
                table_lines.append(f"{key}: {value}")
            elif cells:
                table_lines.append(" | ".join(cells))

    # Combine all text (paragraphs + flattened tables)
    full_text = "\n".join(raw_paragraphs + table_lines)
    full_text = re.sub(r"\s+", " ", full_text)  # normalize all whitespace
    full_text = re.sub(r"\b(\w+)( \1\b)+", r"\1", full_text)  # remove repeated words

    extract_virtual_tags(path, text=full_text)

    return full_text.strip()


# safe_get helper for .msg and.eml to clean stings

def safe_get(obj, key, fallback=""):
    """Safe getter for dicts, objects, and lists with fallback."""
    try:
        if isinstance(obj, dict):
            value = obj.get(key, fallback)
        else:
            value = getattr(obj, key, fallback)

        if value is None:
            return fallback
        if isinstance(value, (list, tuple)):
            return ", ".join(map(str, value))
        return str(value)
    except Exception:
        return fallback


def _extract_msg(path):
    from .fts_core import extract_virtual_tags
    try:
        msg = extract_msg.Message(path)

        subject = safe_get(msg, "subject", "(no subject)")
        sender = safe_get(msg, "sender", "(unknown sender)")
        to = safe_get(msg, "to", "")
        date = safe_get(msg, "date", "")
        # Prioritize plain > RTF > HTML body
        body = safe_get(msg, "body") or safe_get(msg, "bodyRTF") or safe_get(msg, "bodyHTML")

        content = f"Subject: {subject}\nFrom: {sender}\nTo: {to}\nDate: {date}\n{body}"
        content = re.sub(r"\s+", " ", content)

        meta = {"From": sender, "To": to, "Subject": subject, "Date": date}
        extract_virtual_tags(path, meta=meta)

        return content.strip()
    except Exception as e:
        print(f"❌ Failed to extract .msg: {e}")
        return ""


def _extract_eml(path):
    from .fts_core import extract_virtual_tags
    try:
        with open(path, "rb") as f:
            raw_email = f.read()

        ep = eml_parser.EmlParser()
        parsed = ep.decode_email_bytes(raw_email)

        header = parsed.get("header", {})
        subject = safe_get(header, "subject", "(no subject)")
        sender = safe_get(header, "from", ["(unknown sender)"])
        to = safe_get(header, "to", [])
        date = safe_get(header, "date", "(no date)")
        body = safe_get(parsed, "body", [""])

        content = f"Subject: {subject}\nFrom: {sender}\nTo: {to}\nDate: {date}\n{body}"
        content = re.sub(r"\s+", " ", content)

        meta = {"From": sender, "To": to, "Subject": subject, "Date": date}
        extract_virtual_tags(path, meta=meta)

        return content.strip()
    except Exception as e:
        print(f"❌ Failed to extract .eml: {e}")
        return ""



def _extract_pptx(path):
    try:
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        print(f"❌ Failed to extract .pptx: {e}")
        return ""


def _extract_epub(path):
    try:
        book = epub.read_epub(path)
        text = []
        for item in book.get_items():
            if item.get_type() == epub.EpubHtml:
                soup = BeautifulSoup(item.get_body_content(), "html.parser")
                text.append(soup.get_text())
        return "\n".join(text)
    except Exception as e:
        print(f"❌ Failed to extract .epub: {e}")
        return ""


def _extract_odt(path):
    try:
        doc = load(path)
        text = []
        for elem in doc.getElementsByType(P):
            if elem.firstChild:
                text.append(str(elem.firstChild.data))
        return "\n".join(text)
    except Exception as e:
        print(f"❌ Failed to extract .odt: {e}")
        return ""


def _extract_xlsx(path):
    wb = openpyxl.load_workbook(path, data_only=True)
    text = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join([str(cell) for cell in row if cell is not None])
            if row_text.strip():
                text.append(row_text)
    return "\n\n".join(text)


def _extract_pdf(path):
    doc = fitz.open(path)
    full_text = []
    for page in doc:
        text = page.get_text()
        if text.strip():
            full_text.append(text)
        else:
            # OCR fallback
            pix = page.get_pixmap()
            img = Image.open(io.BytesIO(pix.tobytes()))
            ocr_text = pytesseract.image_to_string(img, lang="deu+eng")
            full_text.append(ocr_text)
    return "\n\n".join(full_text)


def _extract_html(path):

    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "html.parser")

    # Remove scripts, styles, hidden navs
    for tag in soup(["script", "style", "nav", "footer", "noscript"]):
        tag.decompose()

    # Title
    title = soup.title.string.strip() if soup.title and soup.title.string else ""

    # Extract headings and paragraphs
    elements = []
    for tag in soup.find_all(["h1", "h2", "h3", "h4", "h5", "p"]):
        txt = tag.get_text(" ", strip=True)
        if txt:
            elements.append(txt)

    combined = f"{title}\n" + "\n".join(elements)
    combined = re.sub(r"\s+", " ", combined)  # collapse whitespace
    return combined.strip()


def store_metadata(path, metadata):
    from .db_utils import connect_db

    if not metadata:
        return
    conn = connect_db()
    cursor = conn.cursor()
    columns = [
        "title",
        "author",
        "subject",
        "created",
        "last_modified",
        "last_modified_by",
        "camera",
        "image_created",
        "dimensions",
        "format",
        "gps",
    ]
    values = [metadata.get(col) for col in columns]
    cursor.execute(
        f"""
        INSERT OR REPLACE INTO file_metadata (path, {', '.join(columns)})
        VALUES (?, {', '.join(['?']*len(columns))})
    """,
        [path] + values,
    )
    conn.commit()
    conn.close()


def _gps_to_decimal(coord, ref):
    try:
        d, m, s = coord

        def to_float(x):
            if hasattr(x, "numerator") and hasattr(x, "denominator"):  # IFDRational
                return float(x.numerator) / float(x.denominator)
            elif isinstance(x, tuple):  # (num, den)
                return float(x[0]) / float(x[1])
            elif isinstance(x, (int, float)):  # already float
                return float(x)
            else:
                print(f"⚠️ Unknown EXIF GPS format: {x} ({type(x)})")
                return 0.0

        val = to_float(d) + to_float(m) / 60.0 + to_float(s) / 3600.0
        return -val if ref in ("S", "W") else val

    except Exception as e:
        print(f"⚠️ Failed to parse GPS coord {coord}: {e}")
        return None


def extract_image_metadata(path: str) -> dict:
    """Extract image metadata including EXIF + GPS if present."""
    md = {}
    try:
        with Image.open(path) as img:
            md["dimensions"] = f"{img.width}x{img.height}"
            md["format"] = img.format

            exif = img._getexif()
            if exif:
                exif_data = {ExifTags.TAGS.get(k, k): v for k, v in exif.items()}
                md["camera"] = exif_data.get("Model") or None
                md["image_created"] = exif_data.get("DateTimeOriginal") or None
                md["title"] = exif_data.get("ImageDescription") or None
                md["author"] = exif_data.get("Artist") or None

                gps = exif_data.get("GPSInfo")
                if isinstance(gps, dict):
                    gps_tags = {ExifTags.GPSTAGS.get(k, k): v for k, v in gps.items()}
                    lat = lon = None
                    if all(
                        k in gps_tags
                        for k in (
                            "GPSLatitude",
                            "GPSLatitudeRef",
                            "GPSLongitude",
                            "GPSLongitudeRef",
                        )
                    ):
                        lat = _gps_to_decimal(
                            gps_tags["GPSLatitude"], gps_tags["GPSLatitudeRef"]
                        )
                        lon = _gps_to_decimal(
                            gps_tags["GPSLongitude"], gps_tags["GPSLongitudeRef"]
                        )
                        md["gps"] = f"{lat:.6f},{lon:.6f}"

        stat = os.stat(path)
        md.setdefault("created", datetime.fromtimestamp(stat.st_ctime).isoformat())
        md.setdefault(
            "last_modified", datetime.fromtimestamp(stat.st_mtime).isoformat()
        )
    except Exception as e:
        print(f"⚠️ Failed to extract image metadata: {path} ({e})")
    return md


def update_file_metadata(file_path, metadata):
    if not metadata:
        return f"Image: {os.path.basename(file_path)}"

    content = f"Image: {os.path.basename(file_path)}"
    for key in [
        "dimensions",
        "format",
        "created",
        "last_modified",
        "camera",
        "image_created",
        "title",
        "author",
        "subject",
        "last_modified_by",
        "gps",
    ]:
        if metadata.get(key):
            content += f" {key}:{metadata[key]}"

    columns = [
        "title",
        "author",
        "subject",
        "created",
        "last_modified",
        "last_modified_by",
        "camera",
        "image_created",
        "dimensions",
        "format",
        "gps",
    ]
    values = [metadata.get(col) for col in columns]

    try:
        conn = sqlite3.connect(DB_FILE)
        cur = conn.cursor()
        cur.execute(
            f"""
            INSERT OR REPLACE INTO file_metadata (path, {', '.join(columns)})
            VALUES (?, {', '.join(['?']*len(columns))})
        """,
            [file_path] + values,
        )
        conn.commit()
        conn.close()
    except Exception as e:
        print(f"⚠️ Failed to update metadata for {file_path}: {e}")

    return content
